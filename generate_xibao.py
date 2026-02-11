"""
财富管理部喜报 & 荣耀战报生成器 - GUI 版
从数据文件（CSV/XLS/XLSX）读取数据，填入 PPTX 模板，
每条数据一页幻灯片，合并输出 PPTX。
带图形界面，适合业务人员使用。
"""

import csv
import copy
import os
import sys
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from collections import defaultdict
from datetime import datetime

from pptx import Presentation


R_NAMESPACE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# 喜报模板占位符映射
FIELD_MAP = {
    "{{分行名称}}": "分行名称",
    "{{客户经理名称}}": "客户经理名称",
    "{{销售额}}": "销售额",
    "{{基金名称}}": "基金产品名称",
}

ROWS_PER_PAGE = 9  # 荣耀战报每页数据行数


def get_base_dir():
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def replace_placeholders_in_paragraph(paragraph, key_map):
    """
    替换段落中的占位符，保留原始格式。
    支持:
      1. 三元组拆分: run "{{", run "key", run "}}"
      2. 合并 run: run "}}-{{"（前一个结尾和后一个开头合并）
      3. 单 run 包含完整占位符: run "{{key}}"
    """
    runs = paragraph.runs
    i = 0
    while i < len(runs):
        # 模式 1: 三元组 run[i]含"{{", run[i+1]=key, run[i+2]含"}}"
        if "{{" in runs[i].text and i + 2 < len(runs):
            potential_key = runs[i + 1].text
            if potential_key in key_map and "}}" in runs[i + 2].text:
                runs[i].text = runs[i].text.replace("{{", "")
                runs[i + 1].text = key_map[potential_key]
                runs[i + 2].text = runs[i + 2].text.replace("}}", "")
                # 用 i+2 而非 i+3，因为 run[i+2] 可能同时含有下一个 "{{"
                i += 2
                continue
        # 模式 2: 单个 run 包含完整占位符 "{{key}}"
        for key, value in key_map.items():
            placeholder = "{{" + key + "}}"
            if placeholder in runs[i].text:
                runs[i].text = runs[i].text.replace(placeholder, value)
        i += 1


def replace_text_in_slide(slide, replacements):
    """替换 slide 中所有文本框的占位符"""
    key_map = {}
    for placeholder, value in replacements.items():
        key = placeholder.strip("{}")
        key_map[key] = value

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                replace_placeholders_in_paragraph(paragraph, key_map)


def replace_text_in_cell(cell, key_map):
    """替换表格单元格中的占位符"""
    for paragraph in cell.text_frame.paragraphs:
        replace_placeholders_in_paragraph(paragraph, key_map)


def copy_slide_from_pptx(target_prs, source_pptx_path):
    src_prs = Presentation(source_pptx_path)
    src_slide = src_prs.slides[0]
    new_slide = target_prs.slides.add_slide(target_prs.slide_layouts[0])
    for child in list(new_slide._element):
        new_slide._element.remove(child)
    rId_map = {}
    for rel in src_slide.part.rels.values():
        if "image" in rel.reltype:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rel.rId] = new_rId
    for child in src_slide._element:
        new_slide._element.append(copy.deepcopy(child))
    for elem in new_slide._element.iter():
        for attr_name in list(elem.attrib.keys()):
            if f"{{{R_NAMESPACE}}}" in attr_name:
                old_rId = elem.attrib[attr_name]
                if old_rId in rId_map:
                    elem.attrib[attr_name] = rId_map[old_rId]


def read_data_file(file_path):
    """读取数据文件，支持 CSV / XLSX / XLS 格式，返回字典列表"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".csv":
        with open(file_path, "r", encoding="utf-8-sig") as f:
            return list(csv.DictReader(f))

    elif ext == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        ws = wb.active
        data = list(ws.iter_rows(values_only=True))
        wb.close()
        if len(data) < 2:
            return []
        headers = [str(h).strip() for h in data[0]]
        return [{headers[j]: (str(cell) if cell is not None else "") for j, cell in enumerate(row)} for row in data[1:] if any(cell is not None for cell in row)]

    elif ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(file_path)
        ws = wb.sheet_by_index(0)
        if ws.nrows < 2:
            return []
        headers = [str(ws.cell_value(0, c)).strip() for c in range(ws.ncols)]
        return [{headers[j]: str(ws.cell_value(r, j)) for j in range(ws.ncols)} for r in range(1, ws.nrows)]

    else:
        raise ValueError(f"不支持的文件格式: {ext}\n请使用 .csv / .xlsx / .xls 文件")


# ===== 喜报生成 =====

def generate_xibao(template_path, data_path, output_path, progress_callback=None):
    """生成喜报（每行数据一页），返回页数"""
    rows = read_data_file(data_path)
    if not rows:
        raise ValueError("数据文件中没有数据")

    total = len(rows)

    # 步骤 1：生成临时 PPTX
    output_dir = os.path.dirname(output_path) or "."
    temp_files = []
    for i, row in enumerate(rows):
        prs = Presentation(template_path)
        slide = prs.slides[0]
        replacements = {ph: row[col] for ph, col in FIELD_MAP.items()}
        replace_text_in_slide(slide, replacements)
        temp_path = os.path.join(output_dir, f"_temp_{i}.pptx")
        prs.save(temp_path)
        temp_files.append(temp_path)
        if progress_callback:
            progress_callback(i + 1, total * 2, f"正在生成第 {i+1}/{total} 页...")

    # 步骤 2：合并
    merged_prs = Presentation(temp_files[0])
    for i, temp_file in enumerate(temp_files[1:], start=1):
        copy_slide_from_pptx(merged_prs, temp_file)
        if progress_callback:
            progress_callback(total + i, total * 2, f"正在合并第 {i+1}/{total} 页...")
    merged_prs.save(output_path)

    # 步骤 3：清理
    for temp_file in temp_files:
        os.remove(temp_file)

    return total


# ===== 荣耀战报生成 =====

def group_data_for_zhanbao(rows):
    """
    按 (分行名称, 基金产品名称) 分组，sum(销售额)，按总额降序排序。
    返回: [{"分行名称": ..., "基金名称": ..., "销售总额": ...}, ...]
    """
    groups = defaultdict(float)
    for row in rows:
        key = (row.get("分行名称", ""), row.get("基金产品名称", ""))
        try:
            amount = float(row.get("销售额", "0").replace(",", ""))
        except ValueError:
            amount = 0
        groups[key] += amount

    result = []
    for (branch, fund), total in groups.items():
        # 格式化金额：保留原始小数精度，去掉尾随零
        total_str = f"{total:g}万"
        result.append({
            "分行名称": branch,
            "基金名称": fund,
            "销售总额": total_str,
        })

    result.sort(key=lambda x: float(x["销售总额"].replace("万", "").replace(",", "")), reverse=True)
    return result


def fill_zhanbao_slide(slide, page_data, start_date, end_date):
    """
    填充荣耀战报的一页 slide。
    - 替换文本框中的日期占位符
    - 替换数据表格中的每行占位符
    - 多余行清空
    """
    # 1. 替换日期文本框
    date_key_map = {
        "数据开始日期": start_date,
        "数据结束日期": end_date,
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                replace_placeholders_in_paragraph(paragraph, date_key_map)

    # 2. 找到数据表格（表格 1，9行数据）并替换
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            table = shape.table
            num_rows = len(list(table.rows))
            # 跳过表头表格（只有1行的是表头）
            if num_rows <= 1:
                continue

            # 这是数据表格
            for ri, row in enumerate(table.rows):
                if ri < len(page_data):
                    # 有数据，填充
                    row_data = page_data[ri]
                    cell_map = {
                        "分行名称": row_data["分行名称"],
                        "基金名称": row_data["基金名称"],
                        "销售总额": row_data["销售总额"],
                    }
                    for ci in range(len(table.columns)):
                        replace_text_in_cell(row.cells[ci], cell_map)
                else:
                    # 没有数据，清空该行
                    for ci in range(len(table.columns)):
                        cell = row.cells[ci]
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                r.text = ""


def get_date_range(rows):
    """
    从数据的"数据日期"列提取最小和最大日期。
    支持常见日期格式：2026/1/2, 2026-1-2, 2026.1.2 等
    返回格式化后的字符串，如 ("1.2", "1.5")
    """
    dates = []
    for row in rows:
        date_str = row.get("数据日期", "").strip()
        if not date_str:
            continue
        # 尝试常见日期格式
        for fmt in ("%Y/%m/%d", "%Y-%m-%d", "%Y.%m.%d",
                     "%Y/%m/%d %H:%M:%S", "%Y-%m-%d %H:%M:%S"):
            try:
                dt = datetime.strptime(date_str.split()[0], fmt.split()[0])
                dates.append(dt)
                break
            except ValueError:
                continue

    if not dates:
        return ("", "")

    min_dt = min(dates)
    max_dt = max(dates)
    # 格式化为 "月.日"
    return (f"{min_dt.month}.{min_dt.day}", f"{max_dt.month}.{max_dt.day}")


def generate_zhanbao(template_path, data_path, output_path,
                     progress_callback=None):
    """
    生成荣耀战报：
    1. 读取数据 -> group by (分行, 基金), sum(销售额), order by desc
    2. 自动从"数据日期"列提取日期区间
    3. 每 9 行一页，填充到 mb2.pptx 模板的表格中
    4. 多页合并
    返回页数
    """
    rows = read_data_file(data_path)
    if not rows:
        raise ValueError("数据文件中没有数据")

    # 自动提取日期区间
    start_date, end_date = get_date_range(rows)

    grouped = group_data_for_zhanbao(rows)
    if not grouped:
        raise ValueError("分组后没有数据")

    # 分页：每页 ROWS_PER_PAGE 行
    pages = []
    for i in range(0, len(grouped), ROWS_PER_PAGE):
        pages.append(grouped[i:i + ROWS_PER_PAGE])

    total_pages = len(pages)
    output_dir = os.path.dirname(output_path) or "."

    # 步骤 1：为每页生成临时 PPTX
    temp_files = []
    for pi, page_data in enumerate(pages):
        prs = Presentation(template_path)
        slide = prs.slides[0]
        fill_zhanbao_slide(slide, page_data, start_date, end_date)
        temp_path = os.path.join(output_dir, f"_temp_zb_{pi}.pptx")
        prs.save(temp_path)
        temp_files.append(temp_path)
        if progress_callback:
            progress_callback(pi + 1, total_pages * 2,
                              f"正在生成第 {pi+1}/{total_pages} 页...")

    # 步骤 2：合并
    merged_prs = Presentation(temp_files[0])
    for i, temp_file in enumerate(temp_files[1:], start=1):
        copy_slide_from_pptx(merged_prs, temp_file)
        if progress_callback:
            progress_callback(total_pages + i, total_pages * 2,
                              f"正在合并第 {i+1}/{total_pages} 页...")
    merged_prs.save(output_path)

    # 步骤 3：清理
    for temp_file in temp_files:
        os.remove(temp_file)

    return total_pages


# ===== GUI =====

class XibaoApp:
    def __init__(self, root):
        self.root = root
        self.root.title("财富管理部喜报生成器")
        self.root.resizable(False, False)

        base_dir = get_base_dir()
        self.default_template = os.path.join(base_dir, "data", "mb.pptx")
        self.default_template_zb = os.path.join(base_dir, "data", "mb2.pptx")
        self.default_data = os.path.join(base_dir, "data", "data.csv")

        self._build_ui()
        self._center_window()

    def _center_window(self):
        self.root.update_idletasks()
        w = self.root.winfo_width()
        h = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (w // 2)
        y = (self.root.winfo_screenheight() // 2) - (h // 2)
        self.root.geometry(f"+{x}+{y}")

    def _build_ui(self):
        # 标题
        title_frame = tk.Frame(self.root, bg="#1a5276", height=60)
        title_frame.pack(fill="x")
        title_frame.pack_propagate(False)
        tk.Label(
            title_frame, text="📋 财富管理部喜报生成器",
            font=("微软雅黑", 16, "bold"), fg="white", bg="#1a5276"
        ).pack(expand=True)

        # 主内容
        main = tk.Frame(self.root, padx=20, pady=15)
        main.pack(fill="both")

        row_idx = 0

        # 数据文件
        tk.Label(main, text="数据文件:", font=("微软雅黑", 10)).grid(
            row=row_idx, column=0, sticky="w", pady=5)
        self.data_var = tk.StringVar(value=self.default_data)
        tk.Entry(main, textvariable=self.data_var, width=50,
                 font=("微软雅黑", 9)).grid(row=row_idx, column=1, padx=5)
        tk.Button(main, text="浏览...",
                  command=self._browse_data).grid(row=row_idx, column=2)
        row_idx += 1

        # 喜报模板文件
        tk.Label(main, text="喜报模板:", font=("微软雅黑", 10)).grid(
            row=row_idx, column=0, sticky="w", pady=5)
        self.template_var = tk.StringVar(value=self.default_template)
        tk.Entry(main, textvariable=self.template_var, width=50,
                 font=("微软雅黑", 9)).grid(row=row_idx, column=1, padx=5)
        tk.Button(main, text="浏览...",
                  command=self._browse_template).grid(row=row_idx, column=2)
        row_idx += 1

        # 战报模板文件
        tk.Label(main, text="战报模板:", font=("微软雅黑", 10)).grid(
            row=row_idx, column=0, sticky="w", pady=5)
        self.template_zb_var = tk.StringVar(value=self.default_template_zb)
        tk.Entry(main, textvariable=self.template_zb_var, width=50,
                 font=("微软雅黑", 9)).grid(row=row_idx, column=1, padx=5)
        tk.Button(main, text="浏览...",
                  command=self._browse_template_zb).grid(row=row_idx, column=2)
        row_idx += 1


        # 进度条
        self.progress = ttk.Progressbar(main, length=450, mode="determinate")
        self.progress.grid(row=row_idx, column=0, columnspan=3,
                           pady=(15, 5), sticky="ew")
        row_idx += 1

        # 状态文字
        self.status_var = tk.StringVar(value="准备就绪")
        tk.Label(main, textvariable=self.status_var,
                 font=("微软雅黑", 9), fg="#666").grid(
            row=row_idx, column=0, columnspan=3, sticky="w")
        row_idx += 1

        # 按钮区
        btn_frame = tk.Frame(main)
        btn_frame.grid(row=row_idx, column=0, columnspan=3, pady=(15, 0))

        self.gen_btn = tk.Button(
            btn_frame, text="🚀  生成喜报", font=("微软雅黑", 11, "bold"),
            bg="#27ae60", fg="white", padx=20, pady=8,
            command=self._on_generate, cursor="hand2"
        )
        self.gen_btn.pack(side="left", padx=8)

        self.gen_zb_btn = tk.Button(
            btn_frame, text="🏆  生成荣耀战报", font=("微软雅黑", 11, "bold"),
            bg="#c0392b", fg="white", padx=20, pady=8,
            command=self._on_generate_zhanbao, cursor="hand2"
        )
        self.gen_zb_btn.pack(side="left", padx=8)

        tk.Button(
            btn_frame, text="📂  打开目录", font=("微软雅黑", 10),
            padx=10, pady=8, command=self._open_output_dir, cursor="hand2"
        ).pack(side="left", padx=8)

    def _browse_template(self):
        path = filedialog.askopenfilename(
            title="选择喜报 PPTX 模板", filetypes=[("PPTX 文件", "*.pptx")])
        if path:
            self.template_var.set(path)

    def _browse_template_zb(self):
        path = filedialog.askopenfilename(
            title="选择战报 PPTX 模板", filetypes=[("PPTX 文件", "*.pptx")])
        if path:
            self.template_zb_var.set(path)

    def _browse_data(self):
        path = filedialog.askopenfilename(
            title="选择数据文件",
            filetypes=[
                ("所有支持的格式", "*.csv *.xlsx *.xls"),
                ("Excel 文件", "*.xlsx *.xls"),
                ("CSV 文件", "*.csv"),
            ]
        )
        if path:
            self.data_var.set(path)

    def _open_output_dir(self):
        output_dir = get_base_dir()
        if os.path.exists(output_dir):
            if sys.platform == "win32":
                os.startfile(output_dir)
            else:
                os.system(f'open "{output_dir}"')
        else:
            messagebox.showinfo("提示", "输出目录不存在")

    def _update_progress(self, current, total, msg):
        self.progress["value"] = (current / total) * 100
        self.status_var.set(msg)
        self.root.update_idletasks()

    def _get_output_name(self, prefix):
        """根据数据文件名生成输出文件名"""
        data_file = self.data_var.get()
        base_name = os.path.splitext(os.path.basename(data_file))[0]
        if base_name == "data":
            base_name = ""
        else:
            base_name = f"_{base_name}"
        return os.path.join(get_base_dir(), f"{prefix}{base_name}.pptx")

    # ----- 喜报 -----

    def _on_generate(self):
        template = self.template_var.get()
        data_file = self.data_var.get()
        output = self._get_output_name("财富管理部喜报")

        if not os.path.exists(template):
            messagebox.showerror("错误", f"模板文件不存在:\n{template}")
            return
        if not os.path.exists(data_file):
            messagebox.showerror("错误", f"数据文件不存在:\n{data_file}")
            return

        self.gen_btn.config(state="disabled", text="⏳  生成中...")
        self.gen_zb_btn.config(state="disabled")
        self.progress["value"] = 0

        def run():
            try:
                count = generate_xibao(
                    template, data_file, output,
                    progress_callback=self._update_progress
                )
                self.root.after(0, lambda: self._on_done(
                    f"✅ 喜报生成完成！共 {count} 页"))
            except Exception as e:
                err_msg = str(e)
                self.root.after(0, lambda: self._on_error(err_msg))

        threading.Thread(target=run, daemon=True).start()

    # ----- 荣耀战报 -----

    def _on_generate_zhanbao(self):
        template = self.template_zb_var.get()
        data_file = self.data_var.get()
        output = self._get_output_name("荣耀战报")

        if not os.path.exists(template):
            messagebox.showerror("错误", f"战报模板文件不存在:\n{template}")
            return
        if not os.path.exists(data_file):
            messagebox.showerror("错误", f"数据文件不存在:\n{data_file}")
            return

        self.gen_zb_btn.config(state="disabled", text="⏳  生成中...")
        self.gen_btn.config(state="disabled")
        self.progress["value"] = 0

        def run():
            try:
                count = generate_zhanbao(
                    template, data_file, output,
                    progress_callback=self._update_progress
                )
                self.root.after(0, lambda: self._on_done(
                    f"✅ 荣耀战报生成完成！共 {count} 页"))
            except Exception as e:
                err_msg = str(e)
                self.root.after(0, lambda: self._on_error(err_msg))

        threading.Thread(target=run, daemon=True).start()

    # ----- 完成 / 错误 -----

    def _on_done(self, msg):
        self.gen_btn.config(state="normal", text="🚀  生成喜报")
        self.gen_zb_btn.config(state="normal", text="🏆  生成荣耀战报")
        self.progress["value"] = 100
        self.status_var.set(msg)

    def _on_error(self, error):
        self.gen_btn.config(state="normal", text="🚀  生成喜报")
        self.gen_zb_btn.config(state="normal", text="🏆  生成荣耀战报")
        self.progress["value"] = 0
        self.status_var.set(f"❌ 生成失败: {error}")
        messagebox.showerror("错误", f"生成失败:\n{error}")


def main():
    root = tk.Tk()
    app = XibaoApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
