"""
综合喜报生成器 - 一键生成完整报告 (包含图片导出功能版)
功能：
1. 读取 data.csv 数据
2. 使用 all.pptx 模板
3. 循环生成每人的喜报并追加到结果 PPT
4. 汇总数据生成战报并追加到结果 PPT 尾部
5. [新增] 自动调用 PowerPoint 将结果导出为图片
"""

import csv
import copy
import os
import sys
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from collections import defaultdict
from datetime import datetime
import subprocess

from pptx import Presentation


R_NAMESPACE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# 喜报占位符
FIELD_MAP = {
    "{{分行名称}}": "分行名称",
    "{{客户经理名称}}": "客户经理名称",
    "{{销售额}}": "销售额",
    "{{基金名称}}": "基金产品名称",
}

# 战报每页行数
ROWS_PER_PAGE_ZHANBAO = 9


def get_base_dir():
    """获取程序运行的基础目录（用于存放生成的输出文件和读取模版）"""
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def replace_placeholders_in_paragraph(paragraph, key_map):
    runs = paragraph.runs
    i = 0
    while i < len(runs):
        if "{{" in runs[i].text and i + 2 < len(runs):
            potential_key = runs[i + 1].text
            if potential_key in key_map and "}}" in runs[i + 2].text:
                runs[i].text = runs[i].text.replace("{{", "")
                runs[i + 1].text = key_map[potential_key]
                runs[i + 2].text = runs[i + 2].text.replace("}}", "")
                i += 2
                continue
        for key, value in key_map.items():
            placeholder = "{{" + key + "}}"
            if placeholder in runs[i].text:
                runs[i].text = runs[i].text.replace(placeholder, value)
        i += 1


def replace_text_in_slide(slide, replacements):
    key_map = {}
    for placeholder, value in replacements.items():
        key = placeholder.strip("{}")
        key_map[key] = value

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                replace_placeholders_in_paragraph(paragraph, key_map)


def replace_text_in_cell(cell, key_map):
    for paragraph in cell.text_frame.paragraphs:
        replace_placeholders_in_paragraph(paragraph, key_map)


def copy_slides_from_pptx(target_prs, source_pptx_path):
    src_prs = Presentation(source_pptx_path)
    for src_slide in src_prs.slides:
        layout_name = src_slide.slide_layout.name
        layout = target_prs.slide_layouts[0]
        for l in target_prs.slide_layouts:
            if l.name == layout_name:
                layout = l
                break
        
        new_slide = target_prs.slides.add_slide(layout)
        
        for child in list(new_slide._element):
            new_slide._element.remove(child)

        rId_map = {}
        for rel in src_slide.part.rels.values():
            if "slideLayout" in rel.reltype:
                continue
            try:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[rel.rId] = new_rId
            except ValueError:
                pass

        for child in src_slide._element:
            new_element = copy.deepcopy(child)
            new_slide._element.append(new_element)
            
            for elem in new_element.iter():
                for attr_name in list(elem.attrib.keys()):
                    if f"{{{R_NAMESPACE}}}" in attr_name:
                        old_rId = elem.attrib[attr_name]
                        if old_rId in rId_map:
                            elem.attrib[attr_name] = rId_map[old_rId]


def read_data_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".csv":
        with open(file_path, "r", encoding="utf-8-sig") as f:
            return list(csv.DictReader(f))

    elif ext == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        ws = wb.active
        data = list(ws.iter_rows(values_only=True))
        wb.close()
        if len(data) < 2:
            return []
        headers = [str(h).strip() for h in data[0]]
        return [{headers[j]: (str(cell) if cell is not None else "") for j, cell in enumerate(row)} for row in data[1:] if any(cell is not None for cell in row)]

    elif ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(file_path)
        ws = wb.sheet_by_index(0)
        if ws.nrows < 2:
            return []
        headers = [str(ws.cell_value(0, c)).strip() for c in range(ws.ncols)]
        return [{headers[j]: str(ws.cell_value(r, j)) for j in range(ws.ncols)} for r in range(1, ws.nrows)]

    else:
        raise ValueError(f"不支持的文件格式: {ext}\n请使用 .csv / .xlsx / .xls 文件")


def get_date_range(rows):
    dates = []
    for row in rows:
        date_str = row.get("数据日期", "").strip()
        if not date_str:
            continue
        for fmt in ("%Y/%m/%d", "%Y-%m-%d", "%Y.%m.%d",
                     "%Y/%m/%d %H:%M:%S", "%Y-%m-%d %H:%M:%S"):
            try:
                dt = datetime.strptime(date_str.split()[0], fmt.split()[0])
                dates.append(dt)
                break
            except ValueError:
                continue

    if not dates:
        return ("", "")

    min_dt = min(dates)
    max_dt = max(dates)
    return (f"{min_dt.month}.{min_dt.day}", f"{max_dt.month}.{max_dt.day}")


def group_data_for_zhanbao(rows):
    groups = defaultdict(float)
    for row in rows:
        key = (row.get("分行名称", ""), row.get("基金产品名称", ""))
        try:
            amount = float(row.get("销售额", "0").replace(",", ""))
        except ValueError:
            amount = 0
        groups[key] += amount

    result = []
    for (branch, fund), total in groups.items():
        total_str = f"{total:g}万" 
        result.append({
            "分行名称": branch,
            "基金名称": fund,
            "销售总额": total_str,
        })

    result.sort(key=lambda x: float(x["销售总额"].replace("万", "").replace(",", "")), reverse=True)
    return result


def fill_zhanbao_slide(slide, page_data, start_date, end_date):
    date_key_map = {
        "数据开始日期": start_date,
        "数据结束日期": end_date,
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                replace_placeholders_in_paragraph(paragraph, date_key_map)

    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            table = shape.table
            num_rows = len(list(table.rows))
            if num_rows <= 1:
                continue

            for ri, row in enumerate(table.rows):
                if ri < len(page_data):
                    row_data = page_data[ri]
                    cell_map = {
                        "分行名称": row_data["分行名称"],
                        "基金名称": row_data["基金名称"],
                        "销售总额": row_data["销售总额"],
                    }
                    for ci in range(len(table.columns)):
                        replace_text_in_cell(row.cells[ci], cell_map)
                else:
                    for ci in range(len(table.columns)):
                        cell = row.cells[ci]
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                r.text = ""


def split_template_by_slides(template_path, temp_dir):
    prs = Presentation(template_path)
    total_slides = len(prs.slides)
    split_files = []

    for i in range(total_slides):
        single_prs = Presentation(template_path)
        xml_slides = single_prs.slides._sldIdLst
        slides = list(xml_slides)
        
        for j, s in enumerate(slides):
            if j != i:
                xml_slides.remove(s)
        
        temp_path = os.path.join(temp_dir, f"_split_template_{i}.pptx")
        single_prs.save(temp_path)
        split_files.append(temp_path)
        
    return split_files


def detect_template_type(prs):
    text_content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content += shape.text_frame.text
            if shape.shape_type == 19:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text_frame.text
    
    if "{{分行名称}}" in text_content or "{{客户经理名称}}" in text_content or "{{基金名称}}" in text_content:
        if "{{数据开始日期}}" in text_content:
            return "SUMMARY"
        return "INDIVIDUAL"
    
    elif "{{数据开始日期}}" in text_content or "{{销售总额}}" in text_content:
        return "SUMMARY"
        
    return "STATIC"


def process_individual_template(template_path, rows, output_dir, index):
    row_temps = []
    for k, row in enumerate(rows):
        p = Presentation(template_path) 
        replacements = {ph: row[col] for ph, col in FIELD_MAP.items()}
        replace_text_in_slide(p.slides[0], replacements)
        t_path = os.path.join(output_dir, f"_part_{index}_row_{k}.pptx")
        p.save(t_path)
        row_temps.append(t_path)

    if not row_temps:
        return None

    merged_prs = Presentation(row_temps[0])
    for t_file in row_temps[1:]:
        copy_slides_from_pptx(merged_prs, t_file)
        try: os.remove(t_file)
        except: pass
    
    try: os.remove(row_temps[0])
    except: pass
        
    part_path = os.path.join(output_dir, f"part_{index}.pptx")
    merged_prs.save(part_path)
    return part_path


def process_summary_template(template_path, rows, output_dir, index):
    start_date, end_date = get_date_range(rows)
    grouped = group_data_for_zhanbao(rows)
    
    if not grouped:
        return None
        
    pages = []
    for i in range(0, len(grouped), ROWS_PER_PAGE_ZHANBAO):
        pages.append(grouped[i:i + ROWS_PER_PAGE_ZHANBAO])
    
    page_temps = []
    for k, page_data in enumerate(pages):
        p = Presentation(template_path)
        fill_zhanbao_slide(p.slides[0], page_data, start_date, end_date)
        t_path = os.path.join(output_dir, f"_part_{index}_page_{k}.pptx")
        p.save(t_path)
        page_temps.append(t_path)
        
    if not page_temps:
        return None

    merged_prs = Presentation(page_temps[0])
    for t_file in page_temps[1:]:
        copy_slides_from_pptx(merged_prs, t_file)
        try: os.remove(t_file)
        except: pass

    try: os.remove(page_temps[0])
    except: pass
        
    part_path = os.path.join(output_dir, f"part_{index}.pptx")
    merged_prs.save(part_path)
    return part_path



# =================================================================
#  完整重写的 generate_full_report (返回 metadata版)
# =================================================================

def generate_full_report(template_path, data_path, output_path, progress_callback=None):
    rows = read_data_file(data_path)
    if not rows:
        raise ValueError("数据文件中没有数据")

    output_dir = os.path.dirname(output_path) or "."
    
    if progress_callback:
        progress_callback(0, 100, "正在分析模板...")
        
    split_templates = split_template_by_slides(template_path, output_dir)
    
    part_info_list = [] # List of (file_path, meta_list)
    
    for i, tmpl_path in enumerate(split_templates):
        prs = Presentation(tmpl_path)
        t_type = detect_template_type(prs)
        
        if progress_callback:
            progress_callback(i * (100 // len(split_templates)), 100, f"正在处理模板页 {i+1} ({t_type})...")
            
        part_file = None
        current_meta = []
        
        if t_type == "INDIVIDUAL":
            part_file = process_individual_template(tmpl_path, rows, output_dir, i)
            # 生成了 len(rows) 张幻灯片
            if part_file:
                 for r_idx in range(len(rows)):
                     current_meta.append({"type": "individual", "row_idx": r_idx})
        
        elif t_type == "SUMMARY":
            # 战报可能有多页，需要计算
            grouped = group_data_for_zhanbao(rows)
            # 计算页数 (虽然 process_summary_template 会处理，但我们需要 meta)
            # 最准确的方法是读取生成的 PPT
            part_file = process_summary_template(tmpl_path, rows, output_dir, i)
            
            if part_file:
                 try:
                     tmp_p = Presentation(part_file)
                     real_count = len(tmp_p.slides)
                     for _ in range(real_count):
                         current_meta.append({"type": "summary"})
                 except:
                     pass

        else:
            part_file = os.path.join(output_dir, f"part_{i}.pptx")
            prs.save(part_file)
            current_meta.append({"type": "static"})
            
        if part_file and os.path.exists(part_file):
            part_info_list.append((part_file, current_meta))
            
        try: os.remove(tmpl_path)
        except: pass

    if progress_callback:
        progress_callback(90, 100, "正在合并所有部分...")
        
    if not part_info_list:
        return 0, []

    # 合并
    first_file, first_meta = part_info_list[0]
    final_prs = Presentation(first_file)
    final_meta = list(first_meta)
    
    for p_file, p_meta in part_info_list[1:]:
        copy_slides_from_pptx(final_prs, p_file)
        final_meta.extend(p_meta)
        
    final_prs.save(output_path)
    
    # 清理分块文件
    try: os.remove(first_file)
    except: pass
    for p_file, _ in part_info_list[1:]:
        try: os.remove(p_file)
        except: pass

    return len(rows), final_meta


# ===== GUI =====

class AllReportsApp:
    def __init__(self, root):
        self.root = root
        self.root.title("财富管理部喜报生成器")
        self.root.geometry("700x580") # 增加高度
        self.root.resizable(False, False)

        base_dir = get_base_dir()
        self.default_template = os.path.join(base_dir, "data", "all.pptx")
        self.default_data = os.path.join(base_dir, "data", "data.csv")
        self.last_output_dir = base_dir

        self._build_ui()
        self._center_window()

    def _center_window(self):
        self.root.update_idletasks()
        w = self.root.winfo_width()
        h = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (w // 2)
        y = (self.root.winfo_screenheight() // 2) - (h // 2)
        self.root.geometry(f"+{x}+{y}")

    def _build_ui(self):
        # --- 配色方案 (Premium Red) ---
        COLOR_PRIMARY = "#B22222"    # Firebrick / Deep Red
        COLOR_ACCENT = "#FFD700"     # Gold
        COLOR_BG = "#F5F5F5"         # White Smoke
        COLOR_CARD = "#FFFFFF"       # Pure White
        COLOR_TEXT = "#333333"       # Dark Gray
        COLOR_TEXT_LIGHT = "#7F8C8D" # Light Gray
        COLOR_BTN_HOVER = "#8B0000"  # Dark Red for hover

        # 字体配置 (跨平台兼容)
        FONT_TITLE = ("Microsoft YaHei UI", 24, "bold") if sys.platform == "win32" else ("PingFang SC", 24, "bold")
        FONT_SUBTITLE = ("Microsoft YaHei UI", 12) if sys.platform == "win32" else ("PingFang SC", 12)
        FONT_LABEL = ("Microsoft YaHei UI", 11) if sys.platform == "win32" else ("PingFang SC", 11)
        FONT_ENTRY = ("Microsoft YaHei UI", 10) if sys.platform == "win32" else ("PingFang SC", 10)
        FONT_BTN_LARGE = ("Microsoft YaHei UI", 14, "bold") if sys.platform == "win32" else ("PingFang SC", 14, "bold")
        FONT_BTN_SMALL = ("Microsoft YaHei UI", 10) if sys.platform == "win32" else ("PingFang SC", 10)

        self.root.configure(bg=COLOR_BG)
        
        self.template_var = tk.StringVar(value=self.default_template)
        self.data_var = tk.StringVar(value=self.default_data)

        # --- 1. 顶部 Header 区域 ---
        header_frame = tk.Frame(self.root, bg=COLOR_PRIMARY, height=120)
        header_frame.pack(fill="x")
        header_frame.pack_propagate(False)

        tk.Label(
            header_frame, text="✨ 财富管理部 · 喜报生成 ✨",
            font=FONT_TITLE, fg=COLOR_ACCENT, bg=COLOR_PRIMARY
        ).pack(expand=True, pady=(20, 0))

        tk.Label(
            header_frame, text="简单、快速、自动化的喜报生成工具",
            font=FONT_SUBTITLE, fg="white", bg=COLOR_PRIMARY
        ).pack(side="bottom", pady=(0, 15))

        # --- 2. 核心内容 Card 区域 ---
        card_frame = tk.Frame(self.root, bg=COLOR_CARD, padx=40, pady=40)
        card_frame.pack(pady=30, padx=30, fill="both", expand=True)
        card_frame.configure(highlightbackground="#E0E0E0", highlightthickness=1)

        def create_input_row(parent, label_text, var, cmd):
            row = tk.Frame(parent, bg=COLOR_CARD)
            row.pack(fill="x", pady=12)
            tk.Label(row, text=label_text, font=FONT_LABEL, fg=COLOR_TEXT, bg=COLOR_CARD, width=22, anchor="w").pack(side="left")
            entry = tk.Entry(row, textvariable=var, font=FONT_ENTRY, bg="#FAFAFA", relief="flat", highlightthickness=1, highlightbackground="#CCCCCC")
            entry.pack(side="left", fill="x", expand=True, ipady=5, padx=10)
            tk.Button(row, text="📂 选择", command=cmd, font=FONT_BTN_SMALL, bg="#EEEEEE", relief="flat", cursor="hand2").pack(side="right", padx=5)

        create_input_row(card_frame, "喜报模板文件 (PPTX)", self.template_var, self._browse_template)
        create_input_row(card_frame, "数据源文件 (CSV/XLSX)", self.data_var, self._browse_data)

        # 进度条
        style = ttk.Style()
        style.theme_use('default')
        style.configure("Red.Horizontal.TProgressbar", background=COLOR_PRIMARY, thickness=10)
        
        self.progress = ttk.Progressbar(card_frame, length=100, mode="determinate", style="Red.Horizontal.TProgressbar")
        self.progress.pack(fill="x", pady=(20, 5))

        # --- 新增：图片导出勾选框 ---
        chk_frame = tk.Frame(card_frame, bg=COLOR_CARD)
        chk_frame.pack(fill="x", pady=(0, 10))
        
        self.export_imgs_var = tk.BooleanVar(value=False)
        chk = tk.Checkbutton(
            chk_frame, 
            text="同时导出为图片 (需要本机安装 PowerPoint/WPS)", 
            variable=self.export_imgs_var,
            font=FONT_BTN_SMALL, bg=COLOR_CARD, fg=COLOR_TEXT,
            selectcolor=COLOR_CARD, activebackground=COLOR_CARD
        )
        chk.pack(side="left")

        # 状态文字
        self.status_var = tk.StringVar(value="准备就绪，等待指令...")
        status_lbl = tk.Label(card_frame, textvariable=self.status_var, font=FONT_LABEL, fg=COLOR_TEXT_LIGHT, bg=COLOR_CARD)
        status_lbl.pack()

        # --- 3. 底部 Action 区域 ---
        action_frame = tk.Frame(self.root, bg=COLOR_BG)
        action_frame.pack(fill="x", pady=20, padx=40)

        self.gen_btn = tk.Button(
            action_frame, text="🔥 生成完整喜报",
            font=FONT_BTN_LARGE, bg=COLOR_PRIMARY, fg="white",
            activebackground=COLOR_BTN_HOVER, activeforeground="white",
            relief="flat", cursor="hand2", padx=20, pady=10,
            command=self._on_generate
        )
        self.gen_btn.pack(side="left", fill="x", expand=True, padx=(0, 10))

        self.open_dir_btn = tk.Button(
            action_frame, text="📂 打开输出文件夹",
            font=FONT_BTN_SMALL, bg="#E0E0E0", fg=COLOR_TEXT,
            relief="flat", cursor="hand2", padx=20, pady=12,
            command=self._open_output_dir
        )
        self.open_dir_btn.pack(side="right", fill="x", padx=(10, 0))

    def _browse_template(self):
        path = filedialog.askopenfilename(filetypes=[("PPTX 文件", "*.pptx")])
        if path: self.template_var.set(path)

    def _browse_data(self):
        path = filedialog.askopenfilename(filetypes=[("数据文件", "*.xlsx *.xls *.csv")])
        if path: self.data_var.set(path)

    def _open_output_dir(self):
        if os.path.exists(self.last_output_dir):
            if sys.platform == "win32":
                os.startfile(self.last_output_dir)
            else:
                os.system(f'open "{self.last_output_dir}"')

    def _on_generate(self):
        template = self.template_var.get()
        data_file = self.data_var.get()
        
        if not os.path.exists(template):
            messagebox.showerror("错误", "模板文件不存在")
            return
        if not os.path.exists(data_file):
            messagebox.showerror("错误", "数据文件不存在")
            return

        try:
            # 预读取一次数据仅用于获取日期命名文件
            rows = read_data_file(data_file)
            if not rows:
                messagebox.showerror("错误", "数据文件中没有数据")
                return
            
            start_date, end_date = get_date_range(rows)
            if start_date and end_date:
                file_name = f"财富管理部喜报({start_date}-{end_date}).pptx"
            else:
                base_name = os.path.splitext(os.path.basename(data_file))[0]
                file_name = f"财富管理部喜报_{base_name}.pptx"
            
            output = os.path.join(get_base_dir(), file_name)
        except Exception as e:
            messagebox.showerror("错误", f"读取数据失败: {e}")
            return

        self.gen_btn.config(state="disabled", text="⏳ 生成中...")
        self.progress["value"] = 0
        self.status_var.set("正在初始化...")
        
        def run():
            try:
                # 调用核心生成逻辑 (现在返回 count, meta)
                count, meta = generate_full_report(template, data_file, output, self._update_progress)
                self.root.after(0, lambda: self._on_ppt_done(output, count, meta))
            except Exception as e:
                err = str(e)
                self.root.after(0, lambda: self._on_error(err))

        threading.Thread(target=run, daemon=True).start()

    def _update_progress(self, current, total, msg):
        self.progress["value"] = (current / total) * 100
        self.status_var.set(msg)
        self.root.update_idletasks()

    def _on_ppt_done(self, output_path, count, meta):
        """PPT 生成完毕，检查是否需要导出图片"""
        if not self.export_imgs_var.get():
            self._finish_all(output_path, count)
            return

        self.status_var.set("📊 正在调用 PowerPoint/WPS 导出图片...")
        self.gen_btn.config(text="⏳ 正在导出图片...")
        
        # 启动转换线程
        threading.Thread(target=self._convert_to_images_thread, args=(output_path, count, meta)).start()

    def _convert_to_images_thread(self, pptx_path, count, meta):
        try:
            base_name_no_ext = os.path.splitext(os.path.basename(pptx_path))[0]
            images_dir = os.path.join(os.path.dirname(pptx_path), f"{base_name_no_ext}_导出图片")
            
            if not os.path.exists(images_dir):
                os.makedirs(images_dir)
            
            # --- 预先构建每一页及其对应的目标文件名 ---
            data_file = self.data_var.get()
            rows = read_data_file(data_file)
            
            # files_map: {slide_index (1-based): absolute_target_path}
            files_map = {}
            
            # 获取日期范围用于战报命名
            start_date, end_date = get_date_range(rows)
            if start_date and end_date:
                zhanbao_base_name = f"战报({start_date}-{end_date})"
            else:
                zhanbao_base_name = "战报"
                
            zhanbao_counter = 1
            
            # 使用 meta 信息构建 map
            for i, m_info in enumerate(meta):
                slide_idx = i + 1
                
                if m_info["type"] == "individual":
                    # 获取对应行数据
                    try:
                        r_idx = m_info["row_idx"]
                        row = rows[r_idx]
                        branch = row.get("分行名称", "未知分行").strip()
                        manager = row.get("客户经理名称", "未知经理").strip()
                        fund = row.get("基金产品名称", "未知产品").strip()
                        safe_name = f"{branch}_{manager}_{fund}".replace("/", "_").replace("\\", "_").replace(":", "")
                        
                        # 如果有重复名字（比如同一个人的两个不同模板生成了2页），需要区分
                        # 这里简单处理：如果 files_map 里已经有了同名的目标路径？
                        # 其实我们的 files_map key 是 slide_idx，value 是 path
                        # 如果 path 重复，覆盖会导致文件覆盖。
                        # 我们检查一下当前目录是否已有同名意图
                        
                        base_fname = f"{safe_name}.jpg"
                        # 检查此文件名是否已被此次任务的其他 slide 占用
                        # 简单起见，如果前面的 slide 已经用了这个名字，我们加个后缀
                        # (虽然通常一人一行数据只对应一组模板，但为了健壮性)
                        target_path = os.path.join(images_dir, base_fname)
                        
                        # 检查已生成的 map 里有没有用过这个 path
                        dup_count = 1
                        while target_path in files_map.values():
                            dup_count += 1
                            target_path = os.path.join(images_dir, f"{safe_name}_{dup_count}.jpg")
                            
                        files_map[slide_idx] = target_path
                    except:
                        pass
                
                elif m_info["type"] == "summary":
                    if zhanbao_counter == 1:
                        fname = f"{zhanbao_base_name}.jpg"
                    else:
                        fname = f"{zhanbao_base_name}_{zhanbao_counter}.jpg"
                    
                    files_map[slide_idx] = os.path.join(images_dir, fname)
                    zhanbao_counter += 1
                
                else:
                    # static or unknown
                    files_map[slide_idx] = os.path.join(images_dir, f"Slide_{slide_idx}.jpg")

            
            # --- 分平台处理 ---
            # 注意：不再需要传递 count 或 zhanbao_base_name 给底层函数，因为 files_map 已经包含了所有信息
            # 但为了兼容之前的 _convert_win32_direct 签名 (它用了 count 和 base_name 来处理 fallback)
            # 我们应该修改 _convert_win32_direct 让它完全依赖 files_map，或者传入 dummy 值
            
            if sys.platform == "win32":
                self._convert_win32_direct(pptx_path, files_map, images_dir)
            else:
                self._convert_mac_workflow(pptx_path, files_map, images_dir)
                
            self.root.after(0, lambda: self._finish_all(pptx_path, count, images_dir))
            
        except Exception as e:
            err = str(e)
            print("Convert Error:", err)
            msg = f"PPT生成成功({count}人)，但导出图片失败。\n需安装Office/WPS。\n错误: {err}"
            self.root.after(0, lambda: messagebox.showwarning("部分完成", msg))
            self.root.after(0, lambda: self._finish_all(pptx_path, count))

    def _convert_win32_direct(self, pptx_path, files_map, output_dir):
        import win32com.client
        
        pptx_path = os.path.abspath(pptx_path)
        
        # 尝试连接 PowerPoint 或 WPS
        try:
            app = win32com.client.Dispatch("PowerPoint.Application")
        except:
            try:
                app = win32com.client.Dispatch("Kwpp.Application")
            except:
                raise Exception("无法调用 PowerPoint 或 WPS，请确认已安装。")
        
        try:
            presentation = app.Presentations.Open(pptx_path, WithWindow=False)
            
            # 获取原始尺寸 (Points)
            sw = presentation.PageSetup.SlideWidth
            sh = presentation.PageSetup.SlideHeight
            
            # 提高分辨率: 设置导出倍数
            # 默认可能是 96DPI，甚至更低。设为 4 倍通常能达到高清效果
            scale = 4
            out_w = int(sw * scale)
            out_h = int(sh * scale)
            
            for i, slide in enumerate(presentation.Slides):
                idx = i + 1
                
                if idx in files_map:
                    target_path = files_map[idx]
                else:
                    # 不再应该发生，因为 files_map 覆盖了 static
                    target_path = os.path.join(output_dir, f"Extra_Slide_{idx}.jpg")
                
                target_path = os.path.abspath(target_path)
                
                try:
                    # 尝试高清导出: slide.Export(FileName, FilterName, ScaleWidth, ScaleHeight)
                    slide.Export(target_path, "JPG", out_w, out_h)
                except:
                    # 如果 WPS 或某些版本不支持宽高参数，回退到默认导出
                    slide.Export(target_path, "JPG")
                
            presentation.Close()
        except Exception as e:
            raise e

    def _convert_mac_workflow(self, pptx_path, files_map, images_dir):
        # Mac 依然使用 AppleScript 全量导出 + 重命名
        temp_dir = os.path.join(images_dir, "temp_export_mac")
        if not os.path.exists(temp_dir): os.makedirs(temp_dir)
        
        import shutil
        for f in os.listdir(temp_dir):
            try: os.remove(os.path.join(temp_dir, f))
            except: pass
            
        self._convert_mac(pptx_path, temp_dir)
        
        for idx in range(1, 9999): 
            found_src = None
            for ext in [".jpg", ".JPG", ".jpeg", ".JPEG", ".png", ".PNG"]:
                t_path = os.path.join(temp_dir, f"Slide{idx}{ext}")
                if os.path.exists(t_path):
                    found_src = t_path
                    break
            
            if not found_src:
                break
            
            if idx in files_map:
                dst_path = files_map[idx]
            else:
                dst_path = os.path.join(images_dir, f"Extra_Slide_{idx}.jpg")
            
            try:
                shutil.move(found_src, dst_path)
            except: pass
            
        try: shutil.rmtree(temp_dir)
        except: pass

    def _finish_all(self, output_path, count, images_dir=None):
        self.gen_btn.config(state="normal", text="🚀 一键生成完整报告")
        self.progress["value"] = 100
        
        msg = f"✅ 完成！生成 {count} 条记录\nPPT: {os.path.basename(output_path)}"
        if images_dir:
            msg += f"\n图片已保存至子文件夹"
            
        self.status_var.set(msg.replace("\n", " "))
        self.last_output_dir = os.path.dirname(output_path) or get_base_dir()
        # messagebox.showinfo("成功", msg) # 之前说不要弹窗

    def _on_error(self, err):
        self.gen_btn.config(state="normal", text="🚀 一键生成完整报告")
        self.progress["value"] = 0
        self.status_var.set(f"出错: {err}")
        messagebox.showerror("失败", f"生成出错:\n{err}")

    def _convert_win32(self, pptx_path, output_dir):
        import win32com.client
        
        pptx_path = os.path.abspath(pptx_path)
        output_dir = os.path.abspath(output_dir)
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # powerpoint.Visible = True 
        try:
            presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
            # ppSaveAsJPG = 17
            presentation.SaveAs(os.path.join(output_dir, "Slide.jpg"), 17)
            presentation.Close()
        finally:
            # powerpoint.Quit() # 慎用Quit
            pass

    def _convert_mac(self, pptx_path, output_dir):
        pptx_path = os.path.abspath(pptx_path)
        output_dir = os.path.abspath(output_dir)
        
        scpt = f'''
        tell application "Microsoft PowerPoint"
            set pptOpen to open "{pptx_path}"
            save pptOpen in "{output_dir}" as save as JPG
            close pptOpen
        end tell
        '''
        p = subprocess.Popen(['osascript', '-e', scpt], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        out, err = p.communicate()
        if p.returncode != 0:
            raise Exception(f"AppleScript Error: {err.decode('utf-8')}")


if __name__ == "__main__":
    root = tk.Tk()
    app = AllReportsApp(root)
    root.mainloop()
