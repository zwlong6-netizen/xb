"""
综合喜报生成器 - 一键生成完整报告
功能：
1. 读取 data.csv 数据
2. 使用 all.pptx 模板 (Slide 0-1 为喜报, Slide 2 为战报)
3. 循环生成每人的喜报并追加到结果 PPT
4. 汇总数据生成战报并追加到结果 PPT 尾部
"""

import csv
import copy
import os
import sys
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from collections import defaultdict
from datetime import datetime

from pptx import Presentation


R_NAMESPACE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# 喜报占位符
FIELD_MAP = {
    "{{分行名称}}": "分行名称",
    "{{客户经理名称}}": "客户经理名称",
    "{{销售额}}": "销售额",
    "{{基金名称}}": "基金产品名称",
}

# 战报每页行数
ROWS_PER_PAGE_ZHANBAO = 9


def get_base_dir():
    """获取程序运行的基础目录（用于存放生成的输出文件和读取模版）"""
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def replace_placeholders_in_paragraph(paragraph, key_map):
    """
    替换段落中的占位符，保留原始格式。
    支持:
      1. 三元组拆分: run "{{", run "key", run "}}"
      2. 合并 run: run "}}-{{"（前一个结尾和后一个开头合并）
      3. 单 run 包含完整占位符: run "{{key}}"
    """
    runs = paragraph.runs
    i = 0
    while i < len(runs):
        # 模式 1: 三元组 run[i]含"{{", run[i+1]=key, run[i+2]含"}}"
        if "{{" in runs[i].text and i + 2 < len(runs):
            potential_key = runs[i + 1].text
            if potential_key in key_map and "}}" in runs[i + 2].text:
                runs[i].text = runs[i].text.replace("{{", "")
                runs[i + 1].text = key_map[potential_key]
                runs[i + 2].text = runs[i + 2].text.replace("}}", "")
                # 用 i+2 而非 i+3，因为 run[i+2] 可能同时含有下一个 "{{"
                i += 2
                continue
        # 模式 2: 单个 run 包含完整占位符 "{{key}}"
        for key, value in key_map.items():
            placeholder = "{{" + key + "}}"
            if placeholder in runs[i].text:
                runs[i].text = runs[i].text.replace(placeholder, value)
        i += 1


def replace_text_in_slide(slide, replacements):
    """替换 slide 中所有文本框的占位符"""
    key_map = {}
    for placeholder, value in replacements.items():
        key = placeholder.strip("{}")
        key_map[key] = value

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                replace_placeholders_in_paragraph(paragraph, key_map)


def replace_text_in_cell(cell, key_map):
    """替换表格单元格中的占位符"""
    for paragraph in cell.text_frame.paragraphs:
        replace_placeholders_in_paragraph(paragraph, key_map)


def copy_slides_from_pptx(target_prs, source_pptx_path):
    """
    从源 PPTX 中复制 **所有** slide 到目标 PPTX。
    """
    src_prs = Presentation(source_pptx_path)
    for src_slide in src_prs.slides:
        # 添加新 slide (尝试匹配 layout，这里简化为第一个 layout 或 blank)
        # 为了更好的兼容性，我们尝试寻找同名 layout
        layout_name = src_slide.slide_layout.name
        layout = target_prs.slide_layouts[0]
        for l in target_prs.slide_layouts:
            if l.name == layout_name:
                layout = l
                break
        
        new_slide = target_prs.slides.add_slide(layout)
        
        # 清空
        for child in list(new_slide._element):
            new_slide._element.remove(child)

        # 复制关系
        rId_map = {}
        for rel in src_slide.part.rels.values():
            if "slideLayout" in rel.reltype:
                continue
            try:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[rel.rId] = new_rId
            except ValueError:
                pass

        # 复制内容
        for child in src_slide._element:
            new_element = copy.deepcopy(child)
            new_slide._element.append(new_element)
            
            for elem in new_element.iter():
                for attr_name in list(elem.attrib.keys()):
                    if f"{{{R_NAMESPACE}}}" in attr_name:
                        old_rId = elem.attrib[attr_name]
                        if old_rId in rId_map:
                            elem.attrib[attr_name] = rId_map[old_rId]


def read_data_file(file_path):
    """读取数据文件，支持 CSV / XLSX / XLS 格式"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".csv":
        with open(file_path, "r", encoding="utf-8-sig") as f:
            return list(csv.DictReader(f))

    elif ext == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        ws = wb.active
        data = list(ws.iter_rows(values_only=True))
        wb.close()
        if len(data) < 2:
            return []
        headers = [str(h).strip() for h in data[0]]
        return [{headers[j]: (str(cell) if cell is not None else "") for j, cell in enumerate(row)} for row in data[1:] if any(cell is not None for cell in row)]

    elif ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(file_path)
        ws = wb.sheet_by_index(0)
        if ws.nrows < 2:
            return []
        headers = [str(ws.cell_value(0, c)).strip() for c in range(ws.ncols)]
        return [{headers[j]: str(ws.cell_value(r, j)) for j in range(ws.ncols)} for r in range(1, ws.nrows)]

    else:
        raise ValueError(f"不支持的文件格式: {ext}\n请使用 .csv / .xlsx / .xls 文件")


def get_date_range(rows):
    """从数据的'数据日期'列提取最小和最大日期"""
    dates = []
    for row in rows:
        date_str = row.get("数据日期", "").strip()
        if not date_str:
            continue
        for fmt in ("%Y/%m/%d", "%Y-%m-%d", "%Y.%m.%d",
                     "%Y/%m/%d %H:%M:%S", "%Y-%m-%d %H:%M:%S"):
            try:
                dt = datetime.strptime(date_str.split()[0], fmt.split()[0])
                dates.append(dt)
                break
            except ValueError:
                continue

    if not dates:
        return ("", "")

    min_dt = min(dates)
    max_dt = max(dates)
    return (f"{min_dt.month}.{min_dt.day}", f"{max_dt.month}.{max_dt.day}")


def group_data_for_zhanbao(rows):
    """汇总数据"""
    groups = defaultdict(float)
    for row in rows:
        key = (row.get("分行名称", ""), row.get("基金产品名称", ""))
        try:
            amount = float(row.get("销售额", "0").replace(",", ""))
        except ValueError:
            amount = 0
        groups[key] += amount

    result = []
    for (branch, fund), total in groups.items():
        total_str = f"{total:g}万" # 保留小数，不取整，加万
        result.append({
            "分行名称": branch,
            "基金名称": fund,
            "销售总额": total_str,
        })

    result.sort(key=lambda x: float(x["销售总额"].replace("万", "").replace(",", "")), reverse=True)
    return result


def fill_zhanbao_slide(slide, page_data, start_date, end_date):
    """填充战报 Slide"""
    # 1. 日期
    date_key_map = {
        "数据开始日期": start_date,
        "数据结束日期": end_date,
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                replace_placeholders_in_paragraph(paragraph, date_key_map)

    # 2. 表格
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            table = shape.table
            num_rows = len(list(table.rows))
            if num_rows <= 1:
                continue

            for ri, row in enumerate(table.rows):
                if ri < len(page_data):
                    row_data = page_data[ri]
                    cell_map = {
                        "分行名称": row_data["分行名称"],
                        "基金名称": row_data["基金名称"],
                        "销售总额": row_data["销售总额"],
                    }
                    for ci in range(len(table.columns)):
                        replace_text_in_cell(row.cells[ci], cell_map)
                else:
                    for ci in range(len(table.columns)):
                        cell = row.cells[ci]
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                r.text = ""


def split_template_by_slides(template_path, temp_dir):
    """
    将模板文件拆分为 N 个单页 PPTX 文件。
    返回: [temp_slide_0.pptx, temp_slide_1.pptx, ...]
    """
    prs = Presentation(template_path)
    total_slides = len(prs.slides)
    split_files = []

    for i in range(total_slides):
        # 加载完整模板，删除除了 i 以外的所有 slide
        single_prs = Presentation(template_path)
        xml_slides = single_prs.slides._sldIdLst
        slides = list(xml_slides)
        
        # 保留 index i，删除其他
        # 注意：需要倒序删除，或者非保留项删除
        for j, s in enumerate(slides):
            if j != i:
                xml_slides.remove(s)
        
        temp_path = os.path.join(temp_dir, f"_split_template_{i}.pptx")
        single_prs.save(temp_path)
        split_files.append(temp_path)
        
    return split_files


def detect_template_type(prs):
    """
    检测模板类型：
    - "INDIVIDUAL": 包含个人字段 ({{分行名称}}, {{客户经理名称}} 等)
    - "SUMMARY": 包含汇总字段 ({{数据开始日期}}, table with {{销售总额}})
    - "STATIC": 其他
    """
    text_content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content += shape.text_frame.text
            if shape.shape_type == 19:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text_frame.text
    
    if "{{分行名称}}" in text_content or "{{客户经理名称}}" in text_content or "{{基金名称}}" in text_content:
        # 如果同时也包含汇总特有字段，优先判定为汇总？
        # 一般个人喜报包含分行/客户经理/基金
        # 汇总包含分行/基金/销售总额列表
        if "{{数据开始日期}}" in text_content:
            return "SUMMARY"
        # 简单判定
        return "INDIVIDUAL"
    
    elif "{{数据开始日期}}" in text_content or "{{销售总额}}" in text_content:
        return "SUMMARY"
        
    return "STATIC"


def process_individual_template(template_path, rows, output_dir, index):
    """处理个人模板：为每行数据生成一页"""
    temp_files = []
    total = len(rows)
    
    # 批量生成，优化速度：
    # 加载模板一次
    # 复制 N 次 slide (内存操作)
    # 填充
    # 保存为一个文件 (part_i.pptx)
    
    prs = Presentation(template_path)
    base_slide = prs.slides[0] # 只有一页
    
    # 复制 base_slide (total - 1) 次
    # 注意：copy_slides_from_pptx 是跨文件，这里是同文件复制
    # 我们用 duplicate_slide (在本文件内复制)
    # 但 duplicate_slide 需要实现。
    # 这里直接用“每次加载模板生成单页存临时文件，最后合并”的笨办法最稳妥
    # 或者：在内存中复制 slide。python-pptx 复制 slide 比较麻烦。
    
    # 采用这方案：
    # 1. 创建 part_i.pptx
    # 2. 循环 rows，每次加载 template_path，填充，append 到 part_i (内存或文件)
    
    # 更高效方案：
    # 加载 template_path -> prs
    # 复制 slide[0] -> slide[1]... slide[N-1]
    # 填充
    # save
    # python-pptx 没有 clone_slide。
    # 所以我们用：生成 N 个单页 temp pptx，然后合并。
    
    # 还是用 generate_xibao 的逻辑：
    # 循环 rows -> 生成 temp_row_k.pptx -> list
    # merge list -> part_i.pptx
    
    row_temps = []
    for k, row in enumerate(rows):
        p = Presentation(template_path) # 只有一页
        replacements = {ph: row[col] for ph, col in FIELD_MAP.items()}
        replace_text_in_slide(p.slides[0], replacements)
        t_path = os.path.join(output_dir, f"_part_{index}_row_{k}.pptx")
        p.save(t_path)
        row_temps.append(t_path)

    # 合并
    if not row_temps:
        return None

    merged_prs = Presentation(row_temps[0])
    for t_file in row_temps[1:]:
        copy_slides_from_pptx(merged_prs, t_file)
        try: os.remove(t_file)
        except: pass
    
    try: os.remove(row_temps[0])
    except: pass
        
    part_path = os.path.join(output_dir, f"part_{index}.pptx")
    merged_prs.save(part_path)
    return part_path


def process_summary_template(template_path, rows, output_dir, index):
    """处理汇总模板"""
    start_date, end_date = get_date_range(rows)
    grouped = group_data_for_zhanbao(rows)
    
    if not grouped:
        return None
        
    pages = []
    for i in range(0, len(grouped), ROWS_PER_PAGE_ZHANBAO):
        pages.append(grouped[i:i + ROWS_PER_PAGE_ZHANBAO])
    
    page_temps = []
    for k, page_data in enumerate(pages):
        p = Presentation(template_path) # 只有一页
        fill_zhanbao_slide(p.slides[0], page_data, start_date, end_date)
        t_path = os.path.join(output_dir, f"_part_{index}_page_{k}.pptx")
        p.save(t_path)
        page_temps.append(t_path)
        
    # 合并
    if not page_temps:
        return None

    merged_prs = Presentation(page_temps[0])
    for t_file in page_temps[1:]:
        copy_slides_from_pptx(merged_prs, t_file)
        try: os.remove(t_file)
        except: pass

    try: os.remove(page_temps[0])
    except: pass
        
    part_path = os.path.join(output_dir, f"part_{index}.pptx")
    merged_prs.save(part_path)
    return part_path


def generate_full_report(template_path, data_path, output_path, progress_callback=None):
    """生成完整报告：基于模板所有 Slide 依次生成"""
    rows = read_data_file(data_path)
    if not rows:
        raise ValueError("数据文件中没有数据")

    output_dir = os.path.dirname(output_path) or "."
    
    # 1. 拆分模板
    if progress_callback:
        progress_callback(0, 100, "正在分析模板...")
        
    split_templates = split_template_by_slides(template_path, output_dir)
    total_steps = len(split_templates) * 2 # process + merge roughly
    current_step = 0
    
    part_files = []
    
    # 2. 对每个拆分后的模板进行处理
    for i, tmpl_path in enumerate(split_templates):
        # 检测类型
        prs = Presentation(tmpl_path)
        t_type = detect_template_type(prs)
        
        if progress_callback:
            progress_callback(i * (100 // len(split_templates)), 100, f"正在处理模板页 {i+1} ({t_type})...")
            
        part_file = None
        if t_type == "INDIVIDUAL":
            part_file = process_individual_template(tmpl_path, rows, output_dir, i)
        elif t_type == "SUMMARY":
            part_file = process_summary_template(tmpl_path, rows, output_dir, i)
        else:
            # STATIC: 直接复制一次作为一页？或者不处理？
            # 假设静态页只保留一份
            part_file = tmpl_path # 直接用原文件作为 part (需要重命名避免混淆吗？)
            # 为了统一，复制一份
            part_file = os.path.join(output_dir, f"part_{i}.pptx")
            prs.save(part_file)
            
        if part_file:
            part_files.append(part_file)
            
        # 清理拆分的模板
        try: os.remove(tmpl_path)
        except: pass

    # 3. 合并所有 Part
    if progress_callback:
        progress_callback(90, 100, "正在合并所有部分...")
        
    if not part_files:
        return 0

    final_prs = Presentation(part_files[0])
    for p_file in part_files[1:]:
        copy_slides_from_pptx(final_prs, p_file)
        
    final_prs.save(output_path)
    
    # 清理 parts
    for p in part_files:
        try: os.remove(p)
        except: pass

    return len(rows)


# ===== GUI =====

class AllReportsApp:
    def __init__(self, root):
        self.root = root
        self.root.title("财富管理部喜报生成器")
        self.root.geometry("600x400")
        self.root.resizable(False, False)

        base_dir = get_base_dir()
        # 默认从程序所在目录的 data 文件夹寻找模板
        self.default_template = os.path.join(base_dir, "data", "all.pptx")
        self.default_data = os.path.join(base_dir, "data", "data.csv")
        self.last_output_dir = base_dir

        self._build_ui()
        self._center_window()

    def _center_window(self):
        self.root.update_idletasks()
        w = self.root.winfo_width()
        h = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (w // 2)
        y = (self.root.winfo_screenheight() // 2) - (h // 2)
        self.root.geometry(f"+{x}+{y}")

    def _build_ui(self):
        # 标题区
        title_frame = tk.Frame(self.root, bg="#2c3e50", height=80)
        title_frame.pack(fill="x")
        title_frame.pack_propagate(False)
        tk.Label(
            title_frame, text="📊 财富管理部喜报生成器",
            font=("微软雅黑", 20, "bold"), fg="white", bg="#2c3e50"
        ).pack(expand=True)
        tk.Label(
            title_frame, text="简单、快速、自动化的报表工具",
            font=("微软雅黑", 10), fg="#bdc3c7", bg="#2c3e50"
        ).pack(side="bottom", pady=5)

        # 内容区
        main = tk.Frame(self.root, padx=30, pady=30)
        main.pack(fill="both", expand=True)

        # 模板选择
        tk.Label(main, text="喜报模版文件（PPTX）:", font=("微软雅黑", 10, "bold")).grid(row=0, column=0, sticky="w", pady=5)
        self.template_var = tk.StringVar(value=self.default_template)
        tk.Entry(main, textvariable=self.template_var, width=40, font=("微软雅黑", 9)).grid(row=0, column=1, padx=5)
        tk.Button(main, text="浏览...", command=self._browse_template).grid(row=0, column=2)

        # 数据选择
        tk.Label(main, text="数据文件 (Excel/CSV):", font=("微软雅黑", 10, "bold")).grid(row=1, column=0, sticky="w", pady=10)
        self.data_var = tk.StringVar(value=self.default_data)
        tk.Entry(main, textvariable=self.data_var, width=40, font=("微软雅黑", 9)).grid(row=1, column=1, padx=5)
        tk.Button(main, text="浏览...", command=self._browse_data).grid(row=1, column=2)

        # 进度条
        self.progress = ttk.Progressbar(main, length=520, mode="determinate")
        self.progress.grid(row=2, column=0, columnspan=3, pady=(20, 5))
        
        self.status_var = tk.StringVar(value="准备就绪")
        tk.Label(main, textvariable=self.status_var, fg="#7f8c8d", font=("微软雅黑", 9)).grid(row=3, column=0, columnspan=3)

        # 按钮
        btn_frame = tk.Frame(main)
        btn_frame.grid(row=4, column=0, columnspan=3, pady=20)

        self.gen_btn = tk.Button(
            btn_frame, text="🚀 一键生成完整报告", font=("微软雅黑", 14, "bold"),
            bg="#27ae60", fg="black", padx=30, pady=10,
            command=self._on_generate, cursor="hand2"
        )
        self.gen_btn.pack(side="left", padx=10)

        self.open_dir_btn = tk.Button(
            btn_frame, text="📁 打开输出目录", font=("微软雅黑", 12, "bold"),
            bg="#34495e", fg="black", padx=20, pady=10,
            command=self._open_output_dir, cursor="hand2"
        )
        self.open_dir_btn.pack(side="left", padx=10)

    def _browse_template(self):
        path = filedialog.askopenfilename(filetypes=[("PPTX 文件", "*.pptx")])
        if path: self.template_var.set(path)

    def _browse_data(self):
        path = filedialog.askopenfilename(filetypes=[("数据文件", "*.xlsx *.xls *.csv")])
        if path: self.data_var.set(path)

    def _open_output_dir(self):
        if os.path.exists(self.last_output_dir):
            if sys.platform == "win32":
                os.startfile(self.last_output_dir)
            else:
                os.system(f'open "{self.last_output_dir}"')

    def _on_generate(self):
        template = self.template_var.get()
        data_file = self.data_var.get()
        
        if not os.path.exists(template):
            messagebox.showerror("错误", "模板文件不存在")
            return
        if not os.path.exists(data_file):
            messagebox.showerror("错误", "数据文件不存在")
            return

        try:
            rows = read_data_file(data_file)
            if not rows:
                messagebox.showerror("错误", "数据文件中没有数据")
                return
            
            start_date, end_date = get_date_range(rows)
            if start_date and end_date:
                file_name = f"财富管理部喜报({start_date}-{end_date}).pptx"
            else:
                base_name = os.path.splitext(os.path.basename(data_file))[0]
                file_name = f"财富管理部喜报_{base_name}.pptx"
            
            output = os.path.join(get_base_dir(), file_name)
        except Exception as e:
            messagebox.showerror("错误", f"读取数据失败: {e}")
            return

        self.gen_btn.config(state="disabled", text="⏳ 生成中...")
        self.progress["value"] = 0
        
        def run():
            try:
                # 传入已经读取好的 rows (可选优化，但目前 generate_full_report 内部还会读一次，暂时保持原样调用)
                count = generate_full_report(template, data_file, output, self._update_progress)
                self.root.after(0, lambda: self._on_done(f"✅ 生成成功！\n文件已保存至: {os.path.basename(output)}", output))
            except Exception as e:
                err = str(e)
                self.root.after(0, lambda: self._on_error(err))

        threading.Thread(target=run, daemon=True).start()

    def _update_progress(self, current, total, msg):
        self.progress["value"] = (current / total) * 100
        self.status_var.set(msg)
        self.root.update_idletasks()

    def _on_done(self, msg, output_path):
        self.gen_btn.config(state="normal", text="🚀 一键生成完整报告")
        self.progress["value"] = 100
        self.status_var.set(f"✅ 完成！文件已保存")
        self.last_output_dir = os.path.dirname(output_path) or get_base_dir()

    def _on_error(self, err):
        self.gen_btn.config(state="normal", text="🚀 一键生成完整报告")
        self.progress["value"] = 0
        self.status_var.set(f"出错: {err}")
        messagebox.showerror("失败", f"生成出错:\n{err}")


if __name__ == "__main__":
    root = tk.Tk()
    app = AllReportsApp(root)
    root.mainloop()
